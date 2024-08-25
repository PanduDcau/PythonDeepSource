import pptx
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Create a presentation object
presentation = pptx.Presentation()

# Define slide layout (1 is title and content layout)
slide_layout = presentation.slide_layouts[1]

# Slide 1 - Title Slide
slide = presentation.slides.add_slide(slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]

title.text = "Blockchain Technology in Supply Chain Management: A Literature Review"
subtitle.text = "An Academic Review and Insights\nPresented by: Group 4 - Taimoor Masud Raja, Threeching Lee\nUniversity of Queensland"

# Slide 2 - Introduction to Blockchain in SCM
slide = presentation.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]

title.text = "Introduction to Blockchain in SCM"
content.text = ("Blockchain Technology: Initially developed for virtual currency, blockchain is now revolutionizing industries, "
                "including Supply Chain Management (SCM).\n\n"
                "Relevance in SCM: Blockchain enhances SCM practices, addressing issues like visibility, product identification, "
                "and flow of goods in global supply chains.\n\n"
                "Key Focus: This review explores blockchain's capacity to improve transparency, traceability, and trust in SCM, "
                "while also acknowledging the challenges in implementation.")

# Slide 3 - Research Questions
slide = presentation.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]

title.text = "Research Questions"
content.text = ("1. How can blockchain technology improve transparency and traceability in supply chain management?\n"
                "2. What are the key challenges and limitations in implementing blockchain solutions in supply chains?\n\n"
                "Significance: Understanding these aspects is crucial for both academic research and practical applications "
                "in the evolving field of SCM.")

# Slide 4 - Summary of Reviewed Articles
slide = presentation.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]

title.text = "Summary of Reviewed Articles"
content.text = ("Overview: The review synthesizes findings from eight recent articles, each exploring different facets of blockchain in SCM.\n\n"
                "Methodologies:\n"
                "- Qualitative Approaches: Interviews and case studies (e.g., Min, 2019; Sivula et al., 2021).\n"
                "- Quantitative Approaches: Surveys and statistical analysis (e.g., Ahmed & MacCarthy, 2022).\n"
                "- Literature Reviews: Systematic reviews of existing research (e.g., Casino et al., 2019).\n\n"
                "Themes: Transparency, traceability, trust, sustainability, and adoption challenges.")

title.text = "References"
content.text = ("Ahmed, W. A. H., & MacCarthy, B. L. (2022). Blockchain in the Supply Chain – A Comprehensive Framework for theory-driven Research. "
                "Digital Business, 2(2), 100043.\n\n"
                "Blossey, G., Eisenhardt, J., & Hahn, G. (2019). Blockchain Technology in Supply Chain Management: An Application Perspective. "
                "Hawaii International Conference on System Sciences.\n\n"
                "Min, H. (2019). Blockchain Technology for Enhancing Supply Chain Resilience. Business Horizons, 62(1), 35–45.\n\n"
                "Dutta, P., Choi, T.-M., Somani, S., & Butala, R. (2020). Blockchain Technology in Supply Chain Operations: Applications, "
                "Challenges, and Research Opportunities. Transportation Research Part E: Logistics and Transportation Review, 142(1), 102067.")

# Save the presentation
presentation_file_path = "/mnt/data/Blockchain_in_SCM_Presentation.pptx"
presentation.save(presentation_file_path)

presentation_file_path